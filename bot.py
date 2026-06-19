<<<<<<< HEAD
import discord
from discord import app_commands, ui
import os
import io
import pypdf
import json
from openai import OpenAI
from datetime import datetime
from dotenv import load_dotenv
 
# --- SETUP ---
import sys

load_dotenv()
TOKEN = os.getenv('DISCORD_TOKEN')
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")    
print("OPENAI-Key geladen:", OPENAI_API_KEY is not None)
client = OpenAI(api_key=os.getenv('OPENAI_API_KEY'))
 
# --- QUIZ LOGIK (BUTTONS) ---
class QuizView(ui.View):
    def __init__(self, questions, user):
        super().__init__(timeout=None)
        self.questions = questions
        self.user = user
        self.current_index = 0
        self.score = 0
        self.add_buttons()
 
    def add_buttons(self):
        self.clear_items()
        q = self.questions[self.current_index]
        for i, option in enumerate(q['options']):
            btn = ui.Button(label=option, custom_id=str(i), style=discord.ButtonStyle.secondary)
            btn.callback = self.check_answer
            self.add_item(btn)
 
    async def check_answer(self, interaction: discord.Interaction):
        if interaction.user != self.user:
            return await interaction.response.send_message("Das ist nicht dein Quiz!", ephemeral=True)
       
        correct_idx = int(self.questions[self.current_index]['correct'])
        if int(interaction.data['custom_id']) == correct_idx:
            self.score += 1
            await interaction.response.send_message("Richtig! ✅", ephemeral=True)
        else:
            await interaction.response.send_message(f"Falsch! ❌ Die richtige Antwort war: {self.questions[self.current_index]['options'][correct_idx]}", ephemeral=True)
 
        self.current_index += 1
        if self.current_index < len(self.questions):
            q = self.questions[self.current_index]
            self.add_buttons()
            await interaction.edit_original_response(content=f"**Frage {self.current_index+1}:** {q['q']}", view=self)
        else:
            # Quiz Ende -> Notenberechnung
            max_p = len(self.questions)
            note = (self.score / max_p) * 5 + 1
            note = round(note * 2) / 2
            await interaction.edit_original_response(content=f"🏁 **Quiz beendet!**\nErgebnis: {self.score}/{max_p}\nDeine Schweizer Note: **{note}**", view=None)
 
# --- BOT HAUPTKLASSE ---
class QuizBot(discord.Client):
    def __init__(self):
        intents = discord.Intents.default()
        intents.message_content = True
        super().__init__(intents=intents)
        self.tree = app_commands.CommandTree(self)
 
    async def setup_hook(self):
        await self.tree.sync()
        print("Bot-System ist vollständig bereit!")
 
bot = QuizBot()
 
# --- BEFEHLE ---
 
@bot.tree.command(name="study", description="Generiert ein Quiz aus Datei, Bild, Screenshot oder Code")
async def study(interaction: discord.Interaction, datei: discord.Attachment):
    await interaction.response.defer()

    try:
        filename = datei.filename.lower()
        file_bytes = await datei.read()
        inhalt = ""
        is_image = filename.endswith((".png", ".jpg", ".jpeg", ".webp"))

        # PDF
        if filename.endswith(".pdf"):
            reader = pypdf.PdfReader(io.BytesIO(file_bytes))
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    inhalt += text + "\n"

        # Word DOCX
        elif filename.endswith(".docx"):
            import docx
            doc = docx.Document(io.BytesIO(file_bytes))
            for paragraph in doc.paragraphs:
                inhalt += paragraph.text + "\n"

        # Excel
        elif filename.endswith((".xlsx", ".xlsm")):
            import openpyxl
            workbook = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
            for sheet in workbook.worksheets:
                inhalt += f"\nTabelle: {sheet.title}\n"
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join([str(cell) for cell in row if cell is not None])
                    if row_text:
                        inhalt += row_text + "\n"

        # PowerPoint
        elif filename.endswith(".pptx"):
            from pptx import Presentation
            presentation = Presentation(io.BytesIO(file_bytes))
            for slide_number, slide in enumerate(presentation.slides, start=1):
                inhalt += f"\nFolie {slide_number}:\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        inhalt += shape.text + "\n"

        # Bilder / Screenshots
        elif is_image:
            inhalt = "Analysiere dieses Bild oder diesen Screenshot und erstelle daraus ein Quiz."

        # Textdateien und Code-Dateien
        elif filename.endswith((
            ".txt", ".md", ".csv",
            ".py", ".java", ".js", ".html", ".css",
            ".sql", ".json", ".xml", ".yml", ".yaml",
            ".php", ".cs", ".cpp", ".c"
        )):
            inhalt = file_bytes.decode("utf-8", errors="ignore")

        else:
            await interaction.followup.send("❌ Dieser Dateityp wird aktuell nicht unterstützt.")
            return

        if not OPENAI_API_KEY:
            await interaction.followup.send("❌ Fehler: Kein OpenAI-Key gefunden. Bitte Key in .env eintragen!")
            return

        # OpenAI Anfrage
        if is_image:
            try:
                models = client.models.list()
                print("MODELS OK")
                print("ERSTES MODELL:", models.data[0].id)
                print("KEY ENDE:", os.getenv("OPENAI_API_KEY")[-4:])
            except Exception as e:
                print("MODELS FEHLER:", e)

            response = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[
                    {
                        "role": "system",
                        "content": """
Erstelle genau 10 Multiple-Choice-Fragen aus dem Bild/Screenshot.
Antworte NUR als valides JSON-Objekt in diesem Format:
{
  "questions": [
    {
      "q": "Frage",
      "options": ["Antwort A", "Antwort B", "Antwort C", "Antwort D"],
      "correct": 0
    }
  ]
}
"""
                    },
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": inhalt},
                            {"type": "image_url", "image_url": {"url": datei.url}}
                        ]
                    }
                ],
                response_format={"type": "json_object"}
            )
        else:
            response = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[
                    {
                        "role": "system",
                        "content": """
Erstelle genau 10 Multiple-Choice-Fragen zum Inhalt.
Antworte NUR als valides JSON-Objekt in diesem Format:
{
  "questions": [
    {
      "q": "Frage",
      "options": ["Antwort A", "Antwort B", "Antwort C", "Antwort D"],
      "correct": 0
    }
  ]
}
"""
                    },
                    {"role": "user", "content": inhalt[:12000]}
                ],
                response_format={"type": "json_object"}
            )

        data = json.loads(response.choices[0].message.content)
        questions = data.get("questions", [])

        if not questions:
            await interaction.followup.send("❌ Es konnten keine Quizfragen erstellt werden.")
            return

        view = QuizView(questions, interaction.user)
        await interaction.followup.send(f"**Frage 1:** {questions[0]['q']}", view=view)

    except Exception as e:
        import traceback
        traceback.print_exc()
        await interaction.followup.send(f"❌ Ein Fehler ist aufgetreten: {type(e).__name__}: {e}")

@bot.tree.command(name="plan", description="Erstellt Lernplan")
async def plan(interaction: discord.Interaction, datum: str):
    try:
        ziel = datetime.strptime(datum, "%d.%m.%Y")
        tage = (ziel - datetime.now()).days
        if tage < 0: return await interaction.response.send_message("Datum in Vergangenheit!")
        minuten = 20 if tage > 14 else 45
        await interaction.response.send_message(f"📅 Prüfung am {datum} ({tage} Tage übrig). Lerne täglich {minuten} Min!")
    except:
        await interaction.response.send_message("Format: DD.MM.YYYY")
 
if __name__ == "__main__":
    bot.run(TOKEN)
 
=======
import discord
from discord import app_commands, ui
import os, io, json, pypdf, logging
from openai import OpenAI
from datetime import datetime, timedelta
from dotenv import load_dotenv

# ─────────────────────────────────────────────
# Setup
# ─────────────────────────────────────────────
load_dotenv()

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    datefmt="%H:%M:%S"
)
log = logging.getLogger("QuizMaster")

TOKEN           = os.getenv("DISCORD_TOKEN")
OPENAI_API_KEY  = os.getenv("OPENAI_API_KEY")

if not TOKEN:
    raise RuntimeError("DISCORD_TOKEN fehlt in .env")
if not OPENAI_API_KEY:
    raise RuntimeError("OPENAI_API_KEY fehlt in .env")

openai_client = OpenAI(api_key=OPENAI_API_KEY)

MAX_FILE_SIZE   = 10 * 1024 * 1024   # 10 MB
MAX_CONTENT     = 12_000             # Zeichen an OpenAI
MAX_QUIZ_LEN    = 14                 # Maximale Lerntage im Plan

# Unterstützte Texterweiterungen
TEXT_EXTENSIONS = (
    ".txt", ".md", ".csv", ".py", ".java", ".js",
    ".html", ".css", ".sql", ".json", ".xml",
    ".yml", ".yaml"
)

# Pro User: {"text": str, "note": float}
user_study_data: dict[int, dict] = {}


# ─────────────────────────────────────────────
# Hilfsfunktionen
# ─────────────────────────────────────────────
def berechne_note(punkte: int, max_punkte: int) -> float:
    """Schweizer Note 1–6, auf 0.5 gerundet."""
    if max_punkte == 0:
        return 1.0
    raw = (punkte * 5 / max_punkte) + 1
    return round(raw * 2) / 2


def lernzeit(note: float) -> int:
    """Lernminuten pro Tag abhängig von der Note."""
    if note >= 5.5:
        return 20
    if note >= 5.0:
        return 25
    if note >= 4.5:
        return 35
    if note >= 4.0:
        return 45
    if note >= 3.0:
        return 60
    return 80


def note_zu_emoji(note: float) -> str:
    if note >= 5.5:
        return "🟢"
    if note >= 4.5:
        return "🟡"
    if note >= 4.0:
        return "🟠"
    return "🔴"


async def send_long_message(interaction: discord.Interaction, text: str):
    """Sendet langen Text in 1900-Zeichen-Blöcken."""
    chunks = [text[i:i + 1900] for i in range(0, len(text), 1900)]
    for chunk in chunks:
        await interaction.followup.send(chunk)


# ─────────────────────────────────────────────
# Datei → Text
# ─────────────────────────────────────────────
async def datei_zu_text(datei: discord.Attachment) -> tuple[str | None, str | None]:
    """Liest eine Discord-Datei und gibt (text, fehler) zurück."""
    filename = datei.filename.lower()

    if datei.size > MAX_FILE_SIZE:
        return None, f"**{datei.filename}** ist zu gross (max. 10 MB)."

    try:
        file_bytes = await datei.read()
    except Exception as e:
        return None, f"**{datei.filename}** konnte nicht gelesen werden: {e}"

    text = ""

    try:
        if filename.endswith(".pdf"):
            reader = pypdf.PdfReader(io.BytesIO(file_bytes))
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"

        elif filename.endswith(".docx"):
            import docx
            doc = docx.Document(io.BytesIO(file_bytes))
            for p in doc.paragraphs:
                if p.text.strip():
                    text += p.text + "\n"

        elif filename.endswith((".xlsx", ".xlsm")):
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
            for sheet in wb.worksheets:
                text += f"\nTabelle: {sheet.title}\n"
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join(str(c) for c in row if c is not None)
                    if row_text.strip():
                        text += row_text + "\n"

        elif filename.endswith(".pptx"):
            from pptx import Presentation
            ppt = Presentation(io.BytesIO(file_bytes))
            for nr, slide in enumerate(ppt.slides, start=1):
                text += f"\nFolie {nr}:\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"

        elif filename.endswith(TEXT_EXTENSIONS):
            text = file_bytes.decode("utf-8", errors="ignore")

        elif filename.endswith((".png", ".jpg", ".jpeg", ".webp")):
            # Bild-URL direkt an OpenAI weitergeben (Vision)
            text = f"[Bild: {datei.url}]"

        else:
            return None, f"**{datei.filename}**: Dateityp wird nicht unterstützt."

    except Exception as e:
        return None, f"**{datei.filename}** konnte nicht verarbeitet werden: {e}"

    if not text.strip():
        return None, f"**{datei.filename}**: Kein lesbarer Inhalt gefunden."

    return text, None


# ─────────────────────────────────────────────
# OpenAI: Quiz erstellen
# ─────────────────────────────────────────────
def openai_quiz_erstellen(inhalt: str) -> list[dict]:
    """Gibt eine Liste von 10 Quizfragen zurück."""
    try:
        response = openai_client.chat.completions.create(
            model="gpt-4o-mini",
            max_tokens=2000,
            messages=[
                {
                    "role": "system",
                    "content": (
                        "Erstelle genau 10 Multiple-Choice-Fragen auf Deutsch zum folgenden Lernstoff.\n"
                        "Jede Frage hat exakt 4 Antwortmöglichkeiten.\n"
                        "Antworte NUR als gültiges JSON ohne Markdown:\n"
                        "{\n"
                        '  "questions": [\n'
                        "    {\n"
                        '      "q": "Frage",\n'
                        '      "options": ["A", "B", "C", "D"],\n'
                        '      "correct": 0\n'
                        "    }\n"
                        "  ]\n"
                        "}"
                    )
                },
                {"role": "user", "content": inhalt[:MAX_CONTENT]}
            ],
            response_format={"type": "json_object"}
        )
        data = json.loads(response.choices[0].message.content)
        questions = data.get("questions", [])
        # Validierung
        valid = []
        for q in questions:
            if (
                isinstance(q, dict)
                and "q" in q
                and "options" in q
                and "correct" in q
                and len(q["options"]) == 4
                and isinstance(q["correct"], int)
                and 0 <= q["correct"] <= 3
            ):
                valid.append(q)
        return valid
    except Exception as e:
        log.error(f"Quiz-Erstellung fehlgeschlagen: {e}")
        return []


# ─────────────────────────────────────────────
# OpenAI: Lernplan erstellen
# ─────────────────────────────────────────────
def openai_lernplan_erstellen(inhalt: str, note: float, datum: str) -> str:
    """Gibt einen formatierten Lernplan-String zurück."""
    heute = datetime.now().date()

    try:
        pruefung = datetime.strptime(datum, "%d.%m.%Y").date()
    except ValueError:
        return "❌ Datumsformat falsch. Nutze z.B. `15.07.2026`."

    if pruefung <= heute:
        return "❌ Das Prüfungsdatum muss in der Zukunft liegen."

    start         = heute + timedelta(days=1)
    letzter_tag   = pruefung - timedelta(days=1)
    lerntage      = max((letzter_tag - start).days + 1, 0)
    minuten       = lernzeit(note)
    emoji         = note_zu_emoji(note)

    if lerntage == 0:
        return (
            f"⚠️ Die Prüfung ist morgen ({datum}).\n"
            f"Wiederhole heute Abend nur die wichtigsten Begriffe – ca. {minuten} Minuten."
        )

    # Themen von OpenAI holen
    try:
        response = openai_client.chat.completions.create(
            model="gpt-4o-mini",
            max_tokens=800,
            messages=[
                {
                    "role": "system",
                    "content": (
                        "Analysiere diesen Lernstoff und extrahiere konkrete Themen auf Deutsch.\n"
                        "Keine Phrasen wie 'lernen' oder 'wiederholen'.\n"
                        "Nutze spezifische Begriffe aus dem Text (z.B. Backup, RAID, AGB, Datenschutz).\n"
                        "Antworte NUR als gültiges JSON ohne Markdown:\n"
                        '{"topics": ["Thema 1", "Thema 2", "Thema 3"]}'
                    )
                },
                {"role": "user", "content": inhalt[:MAX_CONTENT]}
            ],
            response_format={"type": "json_object"}
        )
        data   = json.loads(response.choices[0].message.content)
        topics = data.get("topics", [])
    except Exception as e:
        log.error(f"Themen-Extraktion fehlgeschlagen: {e}")
        topics = []

    if not topics:
        topics = ["Grundlagen wiederholen", "Wichtige Begriffe", "Übungsfragen lösen", "Fehler analysieren"]

    # Plan aufbauen
    tage_im_plan = min(lerntage, MAX_QUIZ_LEN)
    plan_lines = [
        f"📚 **Lernplan**\n",
        f"**Prüfungsdatum:** {datum}",
        f"**Aktuelle Note:** {emoji} **{note}**  _(je schlechter die Note, desto mehr Lernzeit)_",
        f"**Start:** {start.strftime('%d.%m.%Y')}",
        f"**Lernzeit pro Tag:** ca. **{minuten} Minuten**",
        f"**Lerntage bis zur Prüfung:** {lerntage} Tag{'e' if lerntage != 1 else ''}",
    ]

    if lerntage > MAX_QUIZ_LEN:
        plan_lines.append(f"_(Plan zeigt die ersten {MAX_QUIZ_LEN} Tage)_")

    plan_lines.append("")

    for i in range(tage_im_plan):
        tag   = start + timedelta(days=i)
        thema = topics[i % len(topics)]
        plan_lines.append(f"**{tag.strftime('%d.%m.%Y')}:** {thema} – ca. {minuten} Min.")

    plan_lines.append("\n_Das Prüfungsdatum selbst zählt nicht als Lerntag._")

    return "\n".join(plan_lines)


# ─────────────────────────────────────────────
# Quiz-View (Buttons)
# ─────────────────────────────────────────────
class QuizView(ui.View):
    def __init__(self, questions: list[dict], user: discord.User, gesamter_text: str):
        super().__init__(timeout=300)  # 5 Minuten Timeout
        self.questions     = questions
        self.user          = user
        self.gesamter_text = gesamter_text
        self.current_index = 0
        self.score         = 0
        self.results: list[str] = []
        self._render_buttons()

    def _render_buttons(self):
        self.clear_items()
        q = self.questions[self.current_index]
        styles = [
            discord.ButtonStyle.primary,
            discord.ButtonStyle.secondary,
            discord.ButtonStyle.success,
            discord.ButtonStyle.danger,
        ]
        for i, option in enumerate(q["options"]):
            btn = ui.Button(
                label=str(option)[:80],
                custom_id=str(i),
                style=styles[i % len(styles)]
            )
            btn.callback = self._check_answer
            self.add_item(btn)

    async def _check_answer(self, interaction: discord.Interaction):
        if interaction.user.id != self.user.id:
            await interaction.response.send_message("❌ Das ist nicht dein Quiz!", ephemeral=True)
            return

        await interaction.response.defer()

        q           = self.questions[self.current_index]
        correct_idx = int(q["correct"])
        selected    = int(interaction.data["custom_id"])
        num         = self.current_index + 1

        if selected == correct_idx:
            self.score += 1
            self.results.append(f"Frage {num}: _{q['q']}_ → ✅ richtig")
        else:
            richtige_antwort = q["options"][correct_idx]
            self.results.append(
                f"Frage {num}: _{q['q']}_ → ❌ falsch  _(Richtig: {richtige_antwort})_"
            )

        self.current_index += 1

        if self.current_index < len(self.questions):
            self._render_buttons()
            next_q = self.questions[self.current_index]
            fortschritt = f"Frage **{self.current_index + 1}/{len(self.questions)}**"
            await interaction.message.edit(
                content=f"{fortschritt}\n\n{next_q['q']}",
                view=self
            )
        else:
            await self._quiz_beenden(interaction)

    async def _quiz_beenden(self, interaction: discord.Interaction):
        note  = berechne_note(self.score, len(self.questions))
        emoji = note_zu_emoji(note)

        # Study-Daten persistieren
        user_study_data[self.user.id] = {
            "text": self.gesamter_text,
            "note": note
        }
        log.info(f"Quiz beendet: User {self.user} | Note {note} | {self.score}/{len(self.questions)}")

        result_lines = self.results
        result_text  = "\n".join(result_lines)
        if len(result_text) > 1700:
            result_text = result_text[:1700] + "\n_... (gekürzt)_"

        await interaction.message.edit(
            content=(
                f"🏁 **Quiz beendet!**\n\n"
                f"**Ergebnis:** {self.score}/{len(self.questions)}\n"
                f"**Schweizer Note:** {emoji} **{note}**\n\n"
                f"**Auswertung:**\n{result_text}\n\n"
                f"─────────────────────\n"
                f"Nutze jetzt `/plan datum:DD.MM.YYYY` für deinen Lernplan."
            ),
            view=None
        )

    async def on_timeout(self):
        # Buttons deaktivieren wenn Timeout
        for item in self.children:
            item.disabled = True


# ─────────────────────────────────────────────
# Bot
# ─────────────────────────────────────────────
class QuizBot(discord.Client):
    def __init__(self):
        intents = discord.Intents.default()
        super().__init__(intents=intents)
        self.tree = app_commands.CommandTree(self)

    async def setup_hook(self):
        synced = await self.tree.sync()
        log.info(f"{len(synced)} Slash-Commands synchronisiert.")

    async def on_ready(self):
        log.info(f"Bot eingeloggt als {self.user} (ID: {self.user.id})")
        await self.change_presence(
            activity=discord.Activity(
                type=discord.ActivityType.watching,
                name="/study | /plan"
            )
        )


bot = QuizBot()


# ─────────────────────────────────────────────
# /study – Quiz aus Dateien
# ─────────────────────────────────────────────
@bot.tree.command(name="study", description="Erstellt ein Quiz aus deinen Lernunterlagen (bis zu 10 Dateien)")
async def study(
    interaction: discord.Interaction,
    datei1:  discord.Attachment,
    datei2:  discord.Attachment = None,
    datei3:  discord.Attachment = None,
    datei4:  discord.Attachment = None,
    datei5:  discord.Attachment = None,
    datei6:  discord.Attachment = None,
    datei7:  discord.Attachment = None,
    datei8:  discord.Attachment = None,
    datei9:  discord.Attachment = None,
    datei10: discord.Attachment = None,
):
    await interaction.response.defer(thinking=True)

    dateien = [d for d in [
        datei1, datei2, datei3, datei4, datei5,
        datei6, datei7, datei8, datei9, datei10
    ] if d is not None]

    gesamter_text = ""

    for datei in dateien:
        text, fehler = await datei_zu_text(datei)
        if fehler:
            await interaction.followup.send(f"❌ {fehler}")
            return
        gesamter_text += f"\n\n--- {datei.filename} ---\n{text}"

    if not gesamter_text.strip():
        await interaction.followup.send("❌ Kein lesbarer Inhalt in den Dateien gefunden.")
        return

    await interaction.followup.send("⏳ Quizfragen werden erstellt...")

    questions = openai_quiz_erstellen(gesamter_text)

    if not questions:
        await interaction.followup.send("❌ OpenAI konnte keine Quizfragen erstellen. Versuche es erneut.")
        return

    if len(questions) < 5:
        await interaction.followup.send(
            f"⚠️ Nur {len(questions)} gültige Fragen erstellt (statt 10). Das Quiz wird trotzdem gestartet."
        )

    view = QuizView(questions, interaction.user, gesamter_text)
    await interaction.followup.send(
        content=f"Frage **1/{len(questions)}**\n\n{questions[0]['q']}",
        view=view
    )


# ─────────────────────────────────────────────
# /quiz_text – Quiz aus eingetipptem Text
# ─────────────────────────────────────────────
@bot.tree.command(name="quiz_text", description="Erstellt ein Quiz aus direkt eingegebenem Text")
async def quiz_text(interaction: discord.Interaction, text: str):
    await interaction.response.defer(thinking=True)

    if len(text.strip()) < 50:
        await interaction.followup.send("❌ Text ist zu kurz. Gib mindestens 50 Zeichen Lernstoff ein.")
        return

    questions = openai_quiz_erstellen(text)

    if not questions:
        await interaction.followup.send("❌ Keine Quizfragen erstellt. Versuche mehr Text zu geben.")
        return

    view = QuizView(questions, interaction.user, text)
    await interaction.followup.send(
        content=f"Frage **1/{len(questions)}**\n\n{questions[0]['q']}",
        view=view
    )


# ─────────────────────────────────────────────
# /plan – Lernplan aus letztem /study
# ─────────────────────────────────────────────
@bot.tree.command(
    name="plan",
    description="Erstellt deinen Lernplan basierend auf dem letzten /study Quiz"
)
@app_commands.describe(datum="Prüfungsdatum im Format DD.MM.YYYY  (z.B. 15.07.2026)")
async def plan(interaction: discord.Interaction, datum: str):
    await interaction.response.defer(thinking=True)

    daten = user_study_data.get(interaction.user.id)

    if not daten:
        await interaction.followup.send(
            "❌ Kein abgeschlossenes `/study` Quiz gefunden.\n"
            "Mach zuerst ein Quiz mit `/study`, dann kannst du `/plan` nutzen."
        )
        return

    lernplan = openai_lernplan_erstellen(daten["text"], daten["note"], datum)
    await send_long_message(interaction, lernplan)


# ─────────────────────────────────────────────
# /info – Hilfe
# ─────────────────────────────────────────────
@bot.tree.command(name="info", description="Zeigt alle verfügbaren Befehle")
async def info(interaction: discord.Interaction):
    embed = discord.Embed(
        title="📖 QuizMaster – Befehle",
        color=discord.Color.blurple()
    )
    embed.add_field(
        name="/study [datei1] ... [datei10]",
        value="Liest deine Lernunterlagen und startet ein 10-Fragen Quiz.\nUnterstützte Formate: PDF, DOCX, PPTX, XLSX, TXT, MD, CSV, JSON, Bilder",
        inline=False
    )
    embed.add_field(
        name="/quiz_text [text]",
        value="Erstellt ein Quiz aus direkt eingetipptem Text.",
        inline=False
    )
    embed.add_field(
        name="/plan [datum]",
        value="Erstellt einen personalisierten Lernplan basierend auf deinem letzten Quiz.\nFormat: `DD.MM.YYYY`  z.B. `15.07.2026`",
        inline=False
    )
    embed.set_footer(text="Mach zuerst /study, dann /plan – in dieser Reihenfolge.")
    await interaction.response.send_message(embed=embed)


# ─────────────────────────────────────────────
# Start
# ─────────────────────────────────────────────
if __name__ == "__main__":
    bot.run(TOKEN, log_handler=None)
>>>>>>> 5d110c3 (bot.by updated)
